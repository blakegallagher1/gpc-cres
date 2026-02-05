"""
Packaging helpers for one-click exports.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Callable, Dict, List, Optional

from docx import Document as DocxDocument
from openpyxl import Workbook
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

try:
    from pptx import Presentation as PptxPresentation
except ImportError:  # pragma: no cover - optional dependency in some envs
    PptxPresentation = None

Presentation: Optional[Callable[..., Any]] = PptxPresentation


def _ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def _write_pdf(path: Path, title: str, sections: List[str]) -> None:
    c = canvas.Canvas(str(path), pagesize=letter)
    _, height = letter
    y = height - 72
    c.setFont("Helvetica-Bold", 18)
    c.drawString(72, y, title)
    y -= 36
    c.setFont("Helvetica", 11)
    for section in sections:
        for line in section.splitlines():
            if y < 72:
                c.showPage()
                y = height - 72
                c.setFont("Helvetica", 11)
            c.drawString(72, y, line[:120])
            y -= 16
        y -= 8
    c.save()


def generate_investment_memo(
    output_dir: str, project: Dict[str, Any], memo_sections: List[str]
) -> Dict[str, str]:
    out_dir = Path(output_dir)
    _ensure_dir(out_dir)
    doc_path = out_dir / "investment_memo.docx"
    pdf_path = out_dir / "investment_memo.pdf"

    doc = DocxDocument()
    doc.add_heading(project.get("name", "Investment Memo"), level=1)
    for section in memo_sections:
        doc.add_paragraph(section)
    doc.save(str(doc_path))

    _write_pdf(pdf_path, project.get("name", "Investment Memo"), memo_sections)
    return {"docx": str(doc_path), "pdf": str(pdf_path)}


def generate_ic_deck(output_dir: str, project: Dict[str, Any], slides: List[str]) -> Dict[str, str]:
    out_dir = Path(output_dir)
    _ensure_dir(out_dir)
    pptx_path = out_dir / "ic_deck.pptx"
    pdf_path = out_dir / "ic_deck.pdf"

    if Presentation is None:
        raise RuntimeError("python-pptx is required to generate IC decks")
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = project.get("name", "IC Deck")
    slide.placeholders[1].text = project.get("address", "Deal Overview")

    content_layout = prs.slide_layouts[1]
    for section in slides:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = section.splitlines()[0][:60]
        body = slide.shapes.placeholders[1].text_frame
        for line in section.splitlines()[1:]:
            p = body.add_paragraph()
            p.text = line

    prs.save(str(pptx_path))
    _write_pdf(pdf_path, project.get("name", "IC Deck"), slides)
    return {"pptx": str(pptx_path), "pdf": str(pdf_path)}


def generate_underwriting_packet(
    output_dir: str, project: Dict[str, Any], assumptions: Dict[str, Any], results: Dict[str, Any]
) -> Dict[str, str]:
    out_dir = Path(output_dir)
    _ensure_dir(out_dir)
    xlsx_path = out_dir / "underwriting_packet.xlsx"
    pdf_path = out_dir / "underwriting_packet.pdf"

    wb = Workbook()
    ws = wb.active
    ws.title = "Summary"
    ws.append(["Project", project.get("name", "")])
    ws.append(["Address", project.get("address", "")])
    ws.append([])
    ws.append(["Assumptions"])
    for key, value in assumptions.items():
        ws.append([key, value])
    ws.append([])
    ws.append(["Results"])
    for key, value in results.items():
        ws.append([key, value])
    wb.save(str(xlsx_path))

    sections = [
        f"Project: {project.get('name', '')}",
        "Assumptions:\n" + "\n".join([f"{k}: {v}" for k, v in assumptions.items()]),
        "Results:\n" + "\n".join([f"{k}: {v}" for k, v in results.items()]),
    ]
    _write_pdf(pdf_path, "Underwriting Packet", sections)
    return {"xlsx": str(xlsx_path), "pdf": str(pdf_path)}


def generate_dd_report(output_dir: str, project: Dict[str, Any], items: List[Dict[str, Any]]) -> Dict[str, str]:
    out_dir = Path(output_dir)
    _ensure_dir(out_dir)
    doc_path = out_dir / "dd_report.docx"
    pdf_path = out_dir / "dd_report.pdf"

    doc = DocxDocument()
    doc.add_heading(f"DD Report - {project.get('name', '')}", level=1)
    for item in items:
        doc.add_heading(item.get("title", "Checklist Item"), level=2)
        doc.add_paragraph(item.get("summary", ""))
        doc.add_paragraph(f"Status: {item.get('status', 'pending')}")
    doc.save(str(doc_path))

    sections = [
        f"Project: {project.get('name', '')}",
        "\n".join([f"{item.get('title', 'Item')}: {item.get('status', '')}" for item in items]),
    ]
    _write_pdf(pdf_path, "DD Report", sections)
    return {"docx": str(doc_path), "pdf": str(pdf_path)}
